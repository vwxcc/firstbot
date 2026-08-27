import os
import csv
import chardet
from io import BytesIO
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

MAX_TEXT_LENGTH = 20000  # Жесткое ограничение для защиты контекстного окна LLM

def parse_file(file_path: str, filename: str) -> str:
    """
    Фасадная функция для извлечения текста из различных форматов.
    Обеспечивает перехват исключений и предотвращает падение основного процесса.
    """
    ext = os.path.splitext(filename)[1].lower()
    text = ""
    
    try:
        if ext == '.pdf':
            text = _parse_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            text = _parse_docx(file_path)
        elif ext == '.xlsx':
            text = _parse_xlsx(file_path)
        elif ext == '.csv':
            text = _parse_csv(file_path)
        elif ext == '.pptx':
            text = _parse_pptx(file_path)
        elif ext == '.txt':
            text = _parse_txt(file_path)
        else:
            return f"[Системное уведомление: Формат {ext} не поддерживается для прямого извлечения текста]"
            
        if not text or not text.strip():
            return "[Системное уведомление: Файл не содержит текста (возможно, это графический скан без текстового слоя OCR)]"
            
        if len(text) > MAX_TEXT_LENGTH:
            # Усечение текста для предотвращения переполнения токенов
            text = text[:MAX_TEXT_LENGTH] + "\n\n[Системное уведомление: Документ превышает допустимый размер. Текст был усечен для анализа.]"
            
        return text
    except Exception as e:
        print(f"[FILE_PARSER] Ошибка при десериализации {filename}: {e}")
        return "[Системное уведомление: Произошла техническая ошибка при чтении структуры файла]"

def _parse_pdf(file_path: str) -> str:
    text_blocks = []
    with open(file_path, "rb") as f:
        reader = PdfReader(f)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text_blocks.append(extracted)
    return "\n".join(text_blocks)

def _parse_docx(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def _parse_xlsx(file_path: str) -> str:
    # Использование read_only=True минимизирует потребление RAM для больших таблиц
    wb = load_workbook(filename=file_path, data_only=True, read_only=True)
    text_blocks = []
    for sheet in wb.worksheets:
        text_blocks.append(f"--- Таблица/Лист: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            row_data = [str(cell) for cell in row if cell is not None]
            if row_data:
                text_blocks.append(" | ".join(row_data))
    return "\n".join(text_blocks)

def _parse_csv(file_path: str) -> str:
    # Эвристическое определение кодировки файла по первым байтам
    with open(file_path, 'rb') as f:
        raw_data = f.read(10000)
        detected = chardet.detect(raw_data)
        encoding = detected['encoding'] or 'utf-8'
        
    text_blocks = []
    with open(file_path, 'r', encoding=encoding, errors='replace') as f:
        reader = csv.reader(f)
        for row in reader:
            text_blocks.append(" | ".join(row))
    return "\n".join(text_blocks)

def _parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_blocks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_blocks.append(shape.text.strip())
    return "\n".join(text_blocks)

def _parse_txt(file_path: str) -> str:
    with open(file_path, 'rb') as f:
        raw_data = f.read()
        detected = chardet.detect(raw_data)
        encoding = detected['encoding'] or 'utf-8'
    return raw_data.decode(encoding, errors='replace')import os
import csv
import chardet
from io import BytesIO
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

MAX_TEXT_LENGTH = 20000  # Жесткое ограничение для защиты контекстного окна LLM

def parse_file(file_path: str, filename: str) -> str:
    """
    Фасадная функция для извлечения текста из различных форматов.
    Обеспечивает перехват исключений и предотвращает падение основного процесса.
    """
    ext = os.path.splitext(filename)[1].lower()
    text = ""
    
    try:
        if ext == '.pdf':
            text = _parse_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            text = _parse_docx(file_path)
        elif ext == '.xlsx':
            text = _parse_xlsx(file_path)
        elif ext == '.csv':
            text = _parse_csv(file_path)
        elif ext == '.pptx':
            text = _parse_pptx(file_path)
        elif ext == '.txt':
            text = _parse_txt(file_path)
        else:
            return f"[Системное уведомление: Формат {ext} не поддерживается для прямого извлечения текста]"
            
        if not text or not text.strip():
            return "[Системное уведомление: Файл не содержит текста (возможно, это графический скан без текстового слоя OCR)]"
            
        if len(text) > MAX_TEXT_LENGTH:
            # Усечение текста для предотвращения переполнения токенов
            text = text[:MAX_TEXT_LENGTH] + "\n\n[Системное уведомление: Документ превышает допустимый размер. Текст был усечен для анализа.]"
            
        return text
    except Exception as e:
        print(f"[FILE_PARSER] Ошибка при десериализации {filename}: {e}")
        return "[Системное уведомление: Произошла техническая ошибка при чтении структуры файла]"

def _parse_pdf(file_path: str) -> str:
    text_blocks = []
    with open(file_path, "rb") as f:
        reader = PdfReader(f)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text_blocks.append(extracted)
    return "\n".join(text_blocks)

def _parse_docx(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def _parse_xlsx(file_path: str) -> str:
    # Использование read_only=True минимизирует потребление RAM для больших таблиц
    wb = load_workbook(filename=file_path, data_only=True, read_only=True)
    text_blocks = []
    for sheet in wb.worksheets:
        text_blocks.append(f"--- Таблица/Лист: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            row_data = [str(cell) for cell in row if cell is not None]
            if row_data:
                text_blocks.append(" | ".join(row_data))
    return "\n".join(text_blocks)

def _parse_csv(file_path: str) -> str:
    # Эвристическое определение кодировки файла по первым байтам
    with open(file_path, 'rb') as f:
        raw_data = f.read(10000)
        detected = chardet.detect(raw_data)
        encoding = detected['encoding'] or 'utf-8'
        
    text_blocks = []
    with open(file_path, 'r', encoding=encoding, errors='replace') as f:
        reader = csv.reader(f)
        for row in reader:
            text_blocks.append(" | ".join(row))
    return "\n".join(text_blocks)

def _parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_blocks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_blocks.append(shape.text.strip())
    return "\n".join(text_blocks)

def _parse_txt(file_path: str) -> str:
    with open(file_path, 'rb') as f:
        raw_data = f.read()
        detected = chardet.detect(raw_data)
        encoding = detected['encoding'] or 'utf-8'
    return raw_data.decode(encoding, errors='replace')
